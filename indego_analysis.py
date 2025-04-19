import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os
import re
from datetime import datetime

# Constants
DATA_DIR = "data"
OUTPUT_DIR = "output"
TRIPS_FILE = "trips_summary.csv"
SEGMENTS_FILE = "segments_summary.csv"
STATION_FILE = "station_util.csv"
OUTPUT_PPTX = "Perpay_Indego_Presentation.pptx"

def process_raw_data():
    """Process the raw CSV files into the required summary files."""
    # Get all CSV files in the data directory
    csv_files = [f for f in os.listdir(DATA_DIR) if f.endswith('.csv') and not f.startswith('indego-stations-')]
    
    if not csv_files:
        print("No raw CSV files found to process.")
        return
    
    # Initialize dataframes for aggregation
    all_trips = []
    all_segments = []
    station_utilization = {}
    bike_utilization = {}  # New dictionary to store bike utilization data
    
    # Process each CSV file
    for csv_file in csv_files:
        # Read the CSV file
        file_path = os.path.join(DATA_DIR, csv_file)
        try:
            # Read with low_memory=False to avoid DtypeWarning
            df = pd.read_csv(file_path, low_memory=False)
            
            # Print column names for debugging
            print(f"Processing {csv_file} with columns: {df.columns.tolist()}")
            
            # Extract date information
            if 'start_time' in df.columns:
                df['start_time'] = pd.to_datetime(df['start_time'])
                df['start_date'] = df['start_time'].dt.date
                df['start_month'] = df['start_time'].dt.month
                df['start_day'] = df['start_time'].dt.day
                df['start_hour'] = df['start_time'].dt.hour
                
                # Determine year and quarter from the first entry's start_time
                first_date = df['start_time'].min()
                year = first_date.year
                quarter = (first_date.month - 1) // 3 + 1
                year_quarter = f"{year}Q{quarter}"
                
                # Add year_quarter column
                df['year_quarter'] = year_quarter
            
            # Calculate trip duration in minutes if not already present
            if 'duration' not in df.columns and 'start_time' in df.columns and 'end_time' in df.columns:
                df['end_time'] = pd.to_datetime(df['end_time'])
                df['duration'] = (df['end_time'] - df['start_time']).dt.total_seconds() / 60
            
            # Process bike utilization if bike_id and start_station columns exist
            if 'bike_id' in df.columns and 'start_station' in df.columns:
                # System-wide bike utilization
                system_bikes = df['bike_id'].nunique()
                system_trips = len(df)
                
                # Per-station bike utilization
                station_stats = df.groupby('start_station').agg({
                    'bike_id': 'nunique',
                    'trip_id': 'count'
                }).reset_index()
                station_stats.columns = ['station_id', 'unique_bikes', 'total_trips']
                station_stats['trips_per_bike'] = station_stats['total_trips'] / station_stats['unique_bikes']
                
                # Store in bike_utilization dictionary
                if year_quarter not in bike_utilization:
                    bike_utilization[year_quarter] = {
                        'system': {
                            'unique_bikes': system_bikes,
                            'total_trips': system_trips,
                            'trips_per_bike': system_trips / system_bikes
                        },
                        'stations': station_stats.to_dict('records')
                    }
            
            # Append to all_trips
            all_trips.append(df)
            
            # Process segments data
            if 'passholder_type' in df.columns and 'bike_type' in df.columns:
                segments = df.groupby(['year_quarter', 'passholder_type', 'bike_type']).agg({
                    'duration': 'sum',
                    'trip_id': 'count'
                }).reset_index()
                
                segments.columns = ['year_quarter', 'category', 'sub_category', 'total_duration', 'total_trips']
                all_segments.append(segments)
            
            # Process station utilization
            if 'start_station' in df.columns:
                # Count trips per station
                station_counts = df.groupby('start_station').size().reset_index(name='trip_count')
                
                # Initialize station utilization for this quarter if not exists
                if year_quarter not in station_utilization:
                    station_utilization[year_quarter] = {}
                
                # Store station counts
                for _, row in station_counts.iterrows():
                    station_id = row['start_station']
                    if station_id not in station_utilization[year_quarter]:
                        station_utilization[year_quarter][station_id] = 0
                    station_utilization[year_quarter][station_id] += row['trip_count']
            
        except Exception as e:
            print(f"Error processing {csv_file}: {str(e)}")
    
    # Create bike utilization summary file
    if bike_utilization:
        # Create system-wide summary
        system_summary = []
        for year_quarter, data in bike_utilization.items():
            system_summary.append({
                'year_quarter': year_quarter,
                'system_unique_bikes': data['system']['unique_bikes'],
                'system_total_trips': data['system']['total_trips'],
                'system_trips_per_bike': data['system']['trips_per_bike']
            })
        
        system_df = pd.DataFrame(system_summary)
        system_df.to_csv(os.path.join(DATA_DIR, 'bike_util_system.csv'), index=False)
        print("Created bike_util_system.csv")
        
        # Create station-level summary
        station_summary = []
        for year_quarter, data in bike_utilization.items():
            for station in data['stations']:
                station_summary.append({
                    'year_quarter': year_quarter,
                    'station_id': station['station_id'],
                    'unique_bikes': station['unique_bikes'],
                    'total_trips': station['total_trips'],
                    'trips_per_bike': station['trips_per_bike']
                })
        
        station_df = pd.DataFrame(station_summary)
        station_df.to_csv(os.path.join(DATA_DIR, 'bike_util_stations.csv'), index=False)
        print("Created bike_util_stations.csv")
    
    # Combine all data
    if all_trips:
        trips_df = pd.concat(all_trips, ignore_index=True)
        
        # Create trips summary
        trips_summary = trips_df.groupby('year_quarter').agg({
            'trip_id': 'count',
            'duration': ['sum', 'mean'],
            'bike_id': 'nunique'  # Add unique bike count
        }).reset_index()
        
        # Flatten the multi-index columns
        trips_summary.columns = ['year_quarter', 'total_trips', 'total_duration', 'avg_duration', 'unique_bikes']
        
        # Calculate trips per bike
        trips_summary['trips_per_bike'] = trips_summary['total_trips'] / trips_summary['unique_bikes']
        
        # Save to CSV
        trips_summary.to_csv(os.path.join(DATA_DIR, TRIPS_FILE), index=False)
        print(f"Created {TRIPS_FILE}")
    
    if all_segments:
        segments_df = pd.concat(all_segments, ignore_index=True)
        segments_df.to_csv(os.path.join(DATA_DIR, SEGMENTS_FILE), index=False)
        print(f"Created {SEGMENTS_FILE}")
    
    # Create station utilization summary
    if station_utilization:
        # Get all unique station IDs
        all_station_ids = set()
        for quarter_data in station_utilization.values():
            all_station_ids.update(quarter_data.keys())
        
        # Create a dataframe with all stations and quarters
        station_util_data = []
        for station_id in all_station_ids:
            row = {'station_id': station_id}
            for year_quarter in sorted(station_utilization.keys()):
                row[f'util_{year_quarter}'] = station_utilization[year_quarter].get(station_id, 0)
            station_util_data.append(row)
        
        station_util_df = pd.DataFrame(station_util_data)
        station_util_df.to_csv(os.path.join(DATA_DIR, STATION_FILE), index=False)
        print(f"Created {STATION_FILE}")

def load_data():
    """Load and prepare all required data files."""
    trips_df = pd.read_csv(os.path.join(DATA_DIR, TRIPS_FILE))
    segments_df = pd.read_csv(os.path.join(DATA_DIR, SEGMENTS_FILE))
    station_df = pd.read_csv(os.path.join(DATA_DIR, STATION_FILE))
    bike_util_system_df = pd.read_csv(os.path.join(DATA_DIR, 'bike_util_system.csv'))
    bike_util_stations_df = pd.read_csv(os.path.join(DATA_DIR, 'bike_util_stations.csv'))
    
    # Filter out 2016Q4 data
    trips_df = trips_df[trips_df['year_quarter'] != '2016Q4']
    segments_df = segments_df[segments_df['year_quarter'] != '2016Q4']
    bike_util_system_df = bike_util_system_df[bike_util_system_df['year_quarter'] != '2016Q4']
    bike_util_stations_df = bike_util_stations_df[bike_util_stations_df['year_quarter'] != '2016Q4']
    
    return trips_df, segments_df, station_df, bike_util_system_df, bike_util_stations_df

def compute_growth_metrics(trips_df):
    """Compute YoY growth metrics and rolling statistics."""
    # Sort by year_quarter to ensure correct order
    trips_df = trips_df.sort_values('year_quarter')
    
    # Create a copy of the dataframe to avoid SettingWithCopyWarning
    trips_df = trips_df.copy()
    
    # Compute YoY growth for trips and duration
    for metric in ['total_trips', 'total_duration']:
        # Create a mapping of year_quarter to metric value
        year_quarter_map = trips_df.set_index('year_quarter')[metric].to_dict()
        
        # Initialize the yoy_growth column
        growth_col = f'{metric}_yoy_growth'
        trips_df[growth_col] = np.nan
        
        # Calculate YoY growth for each quarter
        for i, row in trips_df.iterrows():
            year_quarter = row['year_quarter']
            year = int(year_quarter.split('Q')[0])
            quarter = int(year_quarter.split('Q')[1])
            
            # Get the previous year's quarter
            prev_year_quarter = f"{year-1}Q{quarter}"
            
            # If we have data for the previous year's quarter, calculate YoY growth
            if prev_year_quarter in year_quarter_map:
                prev_year_value = year_quarter_map[prev_year_quarter]
                current_value = row[metric]
                
                # Calculate YoY growth percentage
                if prev_year_value > 0:
                    trips_df.at[i, growth_col] = ((current_value - prev_year_value) / prev_year_value) * 100
        
        # Calculate long-term geometric average growth rate
        valid_growth_rates = trips_df[growth_col].dropna()
        if not valid_growth_rates.empty:
            # Convert percentage growth rates to multipliers (e.g., 5% -> 1.05)
            growth_multipliers = (valid_growth_rates / 100) + 1
            # Calculate geometric mean and convert back to percentage
            geometric_mean = (np.prod(growth_multipliers) ** (1/len(growth_multipliers)) - 1) * 100
            trips_df[f'{metric}_long_term_avg_growth'] = geometric_mean
        else:
            trips_df[f'{metric}_long_term_avg_growth'] = 0
        
        # Compute rolling statistics (4-quarter window)
        trips_df[f'{metric}_rolling_mean'] = trips_df[growth_col].rolling(window=4, min_periods=1).mean()
        trips_df[f'{metric}_rolling_std'] = trips_df[growth_col].rolling(window=4, min_periods=1).std()
    
    return trips_df

def create_growth_rate_chart(trips_df):
    """Create a chart comparing YoY growth rate with long-term average growth rate."""
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    trips_df = trips_df.sort_values('year_quarter')
    
    # Calculate geometric rolling mean
    growth_multipliers = (trips_df['total_trips_yoy_growth'] / 100) + 1
    
    # Calculate 4-quarter geometric rolling mean
    rolling_geometric_means = []
    for i in range(len(trips_df)):
        if i < 3:  # For the first 3 quarters, use available data
            window = growth_multipliers[0:i+1]
        else:  # For subsequent quarters, use 4-quarter window
            window = growth_multipliers[i-3:i+1]
        
        if len(window) > 0:
            # Calculate geometric mean and convert back to percentage
            geo_mean = (np.prod(window) ** (1/len(window)) - 1) * 100
            rolling_geometric_means.append(geo_mean)
        else:
            rolling_geometric_means.append(np.nan)
    
    trips_df['geometric_rolling_mean'] = rolling_geometric_means
    
    plt.figure(figsize=(16, 9))
    
    # Plot YoY growth
    plt.plot(range(len(trips_df)), trips_df['total_trips_yoy_growth'], 'b-', 
             label='Total Trips YoY Growth', marker='o')
    
    # Plot rolling geometric mean
    plt.plot(range(len(trips_df)), trips_df['geometric_rolling_mean'], 'g--', 
             label='Rolling Geometric Mean')
    
    # Plot long-term average growth rate as a horizontal line
    if 'total_trips_long_term_avg_growth' in trips_df.columns and not trips_df['total_trips_long_term_avg_growth'].isna().all():
        avg_growth = trips_df['total_trips_long_term_avg_growth'].iloc[0]  # All values are the same
        plt.axhline(y=avg_growth, color='r', linestyle='--', 
                   label=f'Long-term Average Growth Rate: {avg_growth:.1f}%')
    
    plt.title('Year-over-Year Total Trips Growth Rate vs. Long-term Average')
    plt.xlabel('Quarter')
    plt.ylabel('Growth Rate (%)')
    
    # Set x-axis ticks to show quarters
    plt.xticks(range(len(trips_df)), trips_df['year_quarter'], rotation=45)
    
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    # Add zero line for reference
    plt.axhline(y=0, color='gray', linestyle='-', alpha=0.3)
    
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_DIR, 'growth_rate.png'))
    plt.close()

def create_duration_growth_chart(trips_df):
    """Create a chart comparing duration YoY growth rate with long-term average growth rate."""
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    trips_df = trips_df.sort_values('year_quarter')
    
    # Calculate geometric rolling mean for duration
    growth_multipliers = (trips_df['total_duration_yoy_growth'] / 100) + 1
    
    # Calculate 4-quarter geometric rolling mean
    rolling_geometric_means = []
    for i in range(len(trips_df)):
        if i < 3:  # For the first 3 quarters, use available data
            window = growth_multipliers[0:i+1]
        else:  # For subsequent quarters, use 4-quarter window
            window = growth_multipliers[i-3:i+1]
        
        if len(window) > 0:
            # Calculate geometric mean and convert back to percentage
            geo_mean = (np.prod(window) ** (1/len(window)) - 1) * 100
            rolling_geometric_means.append(geo_mean)
        else:
            rolling_geometric_means.append(np.nan)
    
    trips_df['duration_geometric_rolling_mean'] = rolling_geometric_means
    
    plt.figure(figsize=(16, 9))
    
    # Plot YoY growth
    plt.plot(range(len(trips_df)), trips_df['total_duration_yoy_growth'], 'b-', 
             label='Duration YoY Growth', marker='o')
    
    # Plot rolling geometric mean
    plt.plot(range(len(trips_df)), trips_df['duration_geometric_rolling_mean'], 'g--', 
             label='Rolling Geometric Mean')
    
    # Plot long-term average growth rate as a horizontal line
    if 'total_duration_long_term_avg_growth' in trips_df.columns and not trips_df['total_duration_long_term_avg_growth'].isna().all():
        avg_growth = trips_df['total_duration_long_term_avg_growth'].iloc[0]  # All values are the same
        plt.axhline(y=avg_growth, color='r', linestyle='--', 
                   label=f'Long-term Average Growth Rate: {avg_growth:.1f}%')
    
    plt.title('Year-over-Year Duration Growth Rate vs. Long-term Average')
    plt.xlabel('Quarter')
    plt.ylabel('Growth Rate (%)')
    
    # Set x-axis ticks to show quarters
    plt.xticks(range(len(trips_df)), trips_df['year_quarter'], rotation=45)
    
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    # Add zero line for reference
    plt.axhline(y=0, color='gray', linestyle='-', alpha=0.3)
    
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_DIR, 'duration_growth_rate.png'))
    plt.close()

def detect_anomalies(trips_df, station_df):
    """Detect anomalous quarters and over-utilized stations."""
    # Flag anomalous quarters based on YoY growth
    trips_df['trips_anomaly'] = abs(trips_df['total_trips_yoy_growth'] - trips_df['total_trips_rolling_mean']) > (2 * trips_df['total_trips_rolling_std'])
    
    # Load station utilization data
    bike_util_stations_df = pd.read_csv(os.path.join(DATA_DIR, 'bike_util_stations.csv'))
    
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    bike_util_stations_df = bike_util_stations_df.sort_values('year_quarter')
    
    # Get the last 4 quarters
    last_4_quarters = sorted(bike_util_stations_df['year_quarter'].unique())[-4:]
    
    # Calculate statistics for each quarter
    quarterly_stats = {}
    for quarter in last_4_quarters:
        quarter_data = bike_util_stations_df[bike_util_stations_df['year_quarter'] == quarter]
        mean_trips = quarter_data['trips_per_bike'].mean()
        std_trips = quarter_data['trips_per_bike'].std()
        quarterly_stats[quarter] = {
            'mean': mean_trips,
            'std': std_trips,
            'data': quarter_data
        }
    
    # Get latest quarter for top 5% calculation
    latest_quarter = last_4_quarters[-1]
    latest_data = quarterly_stats[latest_quarter]['data']
    
    # Calculate top 5% threshold
    top_5_threshold = latest_data['trips_per_bike'].quantile(0.95)
    
    # Initialize dictionary to track stations meeting criteria
    expansion_stations = {}
    
    # Initialize list to store detailed station information
    station_details = []
    
    # Check each station against criteria
    all_stations = bike_util_stations_df['station_id'].unique()
    
    for station in all_stations:
        # Get latest trips per bike for sorting
        latest_station_data = latest_data[latest_data['station_id'] == station]
        latest_trips_per_bike = latest_station_data['trips_per_bike'].iloc[0] if not latest_station_data.empty else 0
        
        # Count deviations in last 4 quarters
        deviations = {1: 0, 2: 0, 3: 0}
        quarters_above = {1: [], 2: [], 3: []}
        deviation_values = {1: [], 2: [], 3: []}  # Store actual deviation values
        
        for quarter in last_4_quarters:
            quarter_data = quarterly_stats[quarter]['data']
            station_data = quarter_data[quarter_data['station_id'] == station]
            
            if not station_data.empty:
                trips = station_data['trips_per_bike'].iloc[0]
                mean = quarterly_stats[quarter]['mean']
                std = quarterly_stats[quarter]['std']
                
                for std_level in [1, 2, 3]:
                    if trips > mean + (std_level * std):
                        deviations[std_level] += 1
                        quarters_above[std_level].append(quarter)
                        deviation_values[std_level].append((trips - mean) / std)  # Store z-score
        
        # Determine the most recent quarter above threshold for each deviation level
        most_recent_quarter = {
            1: max(quarters_above[1]) if quarters_above[1] else '',
            2: max(quarters_above[2]) if quarters_above[2] else '',
            3: max(quarters_above[3]) if quarters_above[3] else ''
        }
        
        # Check criteria and add to details if met
        station_record = None
        
        if deviations[3] >= 1:
            # Calculate average deviation for 3σ
            avg_deviation = np.mean(deviation_values[3]) if deviation_values[3] else 0
            reason = f"Above 3σ in {deviations[3]} quarters (avg {avg_deviation:.1f}σ above mean)"
            expansion_stations[station] = reason
            station_record = {
                'station_id': station,
                'expansion_reason': reason,
                'quarters_above_threshold': ', '.join(quarters_above[3]),
                'deviation_level': 3,
                'quarters_count': deviations[3],
                'most_recent_quarter': most_recent_quarter[3],
                'latest_trips_per_bike': latest_trips_per_bike,
                'sort_priority': 1
            }
        elif deviations[2] >= 2:
            # Calculate average deviation for 2σ
            avg_deviation = np.mean(deviation_values[2]) if deviation_values[2] else 0
            reason = f"Above 2σ in {deviations[2]} quarters (avg {avg_deviation:.1f}σ above mean)"
            expansion_stations[station] = reason
            station_record = {
                'station_id': station,
                'expansion_reason': reason,
                'quarters_above_threshold': ', '.join(quarters_above[2]),
                'deviation_level': 2,
                'quarters_count': deviations[2],
                'most_recent_quarter': most_recent_quarter[2],
                'latest_trips_per_bike': latest_trips_per_bike,
                'sort_priority': 2
            }
        elif deviations[1] >= 3:
            # Calculate average deviation for 1σ
            avg_deviation = np.mean(deviation_values[1]) if deviation_values[1] else 0
            reason = f"Above 1σ in {deviations[1]} quarters (avg {avg_deviation:.1f}σ above mean)"
            expansion_stations[station] = reason
            station_record = {
                'station_id': station,
                'expansion_reason': reason,
                'quarters_above_threshold': ', '.join(quarters_above[1]),
                'deviation_level': 1,
                'quarters_count': deviations[1],
                'most_recent_quarter': most_recent_quarter[1],
                'latest_trips_per_bike': latest_trips_per_bike,
                'sort_priority': 3
            }
        elif not latest_station_data.empty and latest_station_data['trips_per_bike'].iloc[0] >= top_5_threshold:
            # Calculate how far above the threshold
            threshold_exceedance = (latest_trips_per_bike - top_5_threshold) / top_5_threshold * 100
            reason = f"Top 5% in latest quarter ({threshold_exceedance:.1f}% above threshold)"
            expansion_stations[station] = reason
            station_record = {
                'station_id': station,
                'expansion_reason': reason,
                'latest_trips_per_bike': latest_trips_per_bike,
                'threshold_value': top_5_threshold,
                'most_recent_quarter': latest_quarter,
                'sort_priority': 4
            }
        
        if station_record:
            station_details.append(station_record)
    
    # Convert to DataFrame and sort
    df_details = pd.DataFrame(station_details)
    if not df_details.empty:
        df_details = df_details.sort_values(
            by=['latest_trips_per_bike', 'sort_priority', 'most_recent_quarter'],
            ascending=[False, True, False]
        )
        
        # Drop the sorting columns before saving
        df_details = df_details.drop(columns=['sort_priority'])
    
    # Save detailed station information to CSV
    df_details.to_csv(os.path.join(OUTPUT_DIR, 'expansion_recommendations.csv'), index=False)
    
    return trips_df, expansion_stations

def create_trends_chart(trips_df):
    """Create the quarterly trends chart."""
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    trips_df = trips_df.sort_values('year_quarter')
    
    plt.figure(figsize=(16, 9))
    
    # Plot total trips (markers only, no lines)
    ax1 = plt.gca()
    ax1.plot(range(len(trips_df)), trips_df['total_trips'], 'bo', 
             label='Total Trips', markersize=6, linestyle='None')
    
    # Add 1-year (4-quarter) rolling average for total trips
    trips_df['total_trips_rolling_avg'] = trips_df['total_trips'].rolling(window=4, min_periods=1).mean()
    ax1.plot(range(len(trips_df)), trips_df['total_trips_rolling_avg'], 'b:', 
             label='Total Trips (1-Year Rolling Avg)', linewidth=2)
    
    ax1.set_xlabel('Year')
    ax1.set_ylabel('Total Trips', color='b')
    
    # Plot total duration on secondary axis (markers only, no lines)
    ax2 = ax1.twinx()
    ax2.plot(range(len(trips_df)), trips_df['total_duration'], 'ro',
             label='Total Duration', markersize=6, linestyle='None')
    
    # Add 1-year (4-quarter) rolling average for total duration
    trips_df['total_duration_rolling_avg'] = trips_df['total_duration'].rolling(window=4, min_periods=1).mean()
    ax2.plot(range(len(trips_df)), trips_df['total_duration_rolling_avg'], 'r:', 
             label='Total Duration (1-Year Rolling Avg)', linewidth=2)
    
    ax2.set_ylabel('Total Duration (minutes)', color='r')
    
    plt.title('Indego Bike-Share Quarterly Trends')
    
    # Update x-axis to show only years
    years = sorted(set([year_quarter.split('Q')[0] for year_quarter in trips_df['year_quarter']]))
    
    # Calculate positions for year labels (at the middle of each year's quarters)
    year_positions = []
    for year in years:
        year_quarters = [i for i, yq in enumerate(trips_df['year_quarter']) if yq.startswith(year)]
        if year_quarters:
            # Place the year label at the middle of its quarters
            year_positions.append(sum(year_quarters) / len(year_quarters))
    
    # Remove tick labels and set year labels
    plt.xticks(year_positions, years)
    ax1.tick_params(axis='x', rotation=0)  # Set rotation to 0 for better readability
    
    # Add legend
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
    
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_DIR, 'trends.png'))
    plt.close()

def create_anomalies_chart(trips_df):
    """Create the YoY growth control chart."""
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    trips_df = trips_df.sort_values('year_quarter')
    
    # Calculate geometric rolling mean
    # Convert growth rates to multipliers (e.g., 5% -> 1.05)
    growth_multipliers = (trips_df['total_trips_yoy_growth'] / 100) + 1
    
    # Calculate 4-quarter geometric rolling mean
    rolling_geometric_means = []
    for i in range(len(trips_df)):
        if i < 3:  # For the first 3 quarters, use available data
            window = growth_multipliers[0:i+1]
        else:  # For subsequent quarters, use 4-quarter window
            window = growth_multipliers[i-3:i+1]
        
        if len(window) > 0:
            # Calculate geometric mean and convert back to percentage
            geo_mean = (np.prod(window) ** (1/len(window)) - 1) * 100
            rolling_geometric_means.append(geo_mean)
        else:
            rolling_geometric_means.append(np.nan)
    
    trips_df['geometric_rolling_mean'] = rolling_geometric_means
    
    plt.figure(figsize=(16, 9))
    
    # Plot YoY growth
    plt.plot(range(len(trips_df)), trips_df['total_trips_yoy_growth'], 'b-', 
             label='Total Trips YoY Growth', marker='o')
    
    # Plot control bands
    plt.fill_between(range(len(trips_df)),
                    trips_df['geometric_rolling_mean'] + 2 * trips_df['total_trips_rolling_std'],
                    trips_df['geometric_rolling_mean'] - 2 * trips_df['total_trips_rolling_std'],
                    alpha=0.2, color='gray', label='±2σ Band')
    
    plt.plot(range(len(trips_df)), trips_df['geometric_rolling_mean'], 'r--', 
             label='Rolling Geometric Mean')
    
    plt.title('Year-over-Year Total Trips Growth Control Chart')
    plt.xlabel('Quarter')
    plt.ylabel('YoY Growth (%)')
    
    # Set x-axis ticks to show quarters
    plt.xticks(range(len(trips_df)), trips_df['year_quarter'], rotation=45)
    
    plt.legend()
    plt.grid(True, alpha=0.3)
    
    # Add zero line for reference
    plt.axhline(y=0, color='gray', linestyle='-', alpha=0.3)
    
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_DIR, 'anomalies.png'))
    plt.close()

def create_presentation(trips_df, expansion_stations):
    """Create the PowerPoint presentation."""
    prs = Presentation()
    
    # Slide 1: Growth Overview
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide1.shapes.title
    title.text = "Program Growth Overview"
    
    # Add trends chart
    left = Inches(0.5)
    top = Inches(1.5)
    pic_width = Inches(9)
    pic_height = Inches(5)
    slide1.shapes.add_picture(os.path.join(OUTPUT_DIR, 'trends.png'), left, top, width=pic_width, height=pic_height)
    
    # Add bullet points
    left = Inches(0.5)
    top = Inches(7)
    width = Inches(9)
    height = Inches(1)
    
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    
    # Latest metrics
    latest = trips_df.iloc[-1]
    p = tf.add_paragraph()
    p.text = f"• Latest YoY growth: {latest['total_trips_yoy_growth']:.1f}% (Baseline: {latest['total_trips_rolling_mean']:.1f}%)"
    
    p = tf.add_paragraph()
    p.text = f"• Long-term average growth rate: {latest['total_trips_long_term_avg_growth']:.1f}%"
    
    p = tf.add_paragraph()
    p.text = f"• Average trip duration trend: {'Increasing' if latest['avg_duration'] > trips_df['avg_duration'].mean() else 'Decreasing'}"
    
    # Slide 2: Growth Rate Analysis
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide2.shapes.title
    title.text = "Growth Rate Analysis"
    
    # Add growth rate chart
    slide2.shapes.add_picture(os.path.join(OUTPUT_DIR, 'growth_rate.png'), left, top, width=pic_width, height=pic_height)
    
    # Add bullet points
    textbox = slide2.shapes.add_textbox(left, Inches(7), width, height)
    tf = textbox.text_frame
    
    # Growth rate insights
    p = tf.add_paragraph()
    p.text = f"• Current growth rate: {latest['total_trips_yoy_growth']:.1f}% vs. long-term average: {latest['total_trips_long_term_avg_growth']:.1f}%"
    
    p = tf.add_paragraph()
    p.text = f"• Growth trend: {'Accelerating' if latest['total_trips_yoy_growth'] > latest['total_trips_rolling_mean'] else 'Decelerating'}"
    
    # Slide 3: Duration Growth Analysis
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide3.shapes.title
    title.text = "Duration Growth Analysis"
    
    # Add duration growth chart
    slide3.shapes.add_picture(os.path.join(OUTPUT_DIR, 'duration_growth_rate.png'), left, top, width=pic_width, height=pic_height)
    
    # Add bullet points
    textbox = slide3.shapes.add_textbox(left, Inches(7), width, height)
    tf = textbox.text_frame
    
    # Duration growth insights
    p = tf.add_paragraph()
    p.text = f"• Current duration growth rate: {latest['total_duration_yoy_growth']:.1f}% vs. long-term average: {latest['total_duration_long_term_avg_growth']:.1f}%"
    
    p = tf.add_paragraph()
    p.text = f"• Duration growth trend: {'Accelerating' if latest['total_duration_yoy_growth'] > latest['total_duration_rolling_mean'] else 'Decelerating'}"
    
    # Slide 4: Anomalies & Recommendations
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide4.shapes.title
    title.text = "Anomaly Detection & Recommendations"
    
    # Add anomalies chart
    slide4.shapes.add_picture(os.path.join(OUTPUT_DIR, 'anomalies.png'), left, top, width=pic_width, height=pic_height)
    
    # Add bullet points
    textbox = slide4.shapes.add_textbox(left, Inches(7), width, height)
    tf = textbox.text_frame
    
    # Anomalous quarters
    anomalous_quarters = trips_df[trips_df['trips_anomaly']]['year_quarter'].tolist()
    p = tf.add_paragraph()
    p.text = f"• Anomalous quarters: {', '.join(anomalous_quarters) if anomalous_quarters else 'None detected'}"
    
    # Station expansion recommendation (simplified)
    p = tf.add_paragraph()
    p.text = f"• Recommend station expansion in {len(expansion_stations)} locations (see expansion_recommendations.csv for details)"
    
    # Save presentation
    prs.save(os.path.join(OUTPUT_DIR, OUTPUT_PPTX))

def create_bike_utilization_chart(bike_util_system_df):
    """Create a chart showing system-wide bike utilization trends."""
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    bike_util_system_df = bike_util_system_df.sort_values('year_quarter')
    
    plt.figure(figsize=(16, 9))
    
    # Plot unique bikes and trips per bike
    ax1 = plt.gca()
    ax1.plot(range(len(bike_util_system_df)), bike_util_system_df['system_unique_bikes'], 
             'b-', label='Unique Bikes', marker='o')
    ax1.set_xlabel('Year')
    ax1.set_ylabel('Number of Unique Bikes', color='b')
    
    # Plot trips per bike on secondary axis
    ax2 = ax1.twinx()
    ax2.plot(range(len(bike_util_system_df)), bike_util_system_df['system_trips_per_bike'], 
             'r-', label='Trips per Bike', marker='o')
    ax2.set_ylabel('Trips per Bike', color='r')
    
    plt.title('System-wide Bike Utilization Trends')
    
    # Update x-axis to show only years
    years = sorted(set([year_quarter.split('Q')[0] for year_quarter in bike_util_system_df['year_quarter']]))
    
    # Calculate positions for year labels (at the middle of each year's quarters)
    year_positions = []
    for year in years:
        year_quarters = [i for i, yq in enumerate(bike_util_system_df['year_quarter']) if yq.startswith(year)]
        if year_quarters:
            # Place the year label at the middle of its quarters
            year_positions.append(sum(year_quarters) / len(year_quarters))
    
    # Remove tick labels and set year labels
    plt.xticks(year_positions, years)
    ax1.tick_params(axis='x', rotation=0)  # Set rotation to 0 for better readability
    
    # Add legend
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
    
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_DIR, 'bike_utilization.png'))
    plt.close()

def create_station_bike_utilization_chart(bike_util_stations_df):
    """Create a chart showing station-level bike utilization trends."""
    # Sort by year_quarter to ensure correct ordering (oldest to newest)
    bike_util_stations_df = bike_util_stations_df.sort_values('year_quarter')
    
    plt.figure(figsize=(16, 9))
    
    # Calculate statistics per quarter
    quarterly_stats = bike_util_stations_df.groupby('year_quarter').agg({
        'unique_bikes': ['mean', 'std'],
        'trips_per_bike': ['mean', 'std']
    }).reset_index()
    
    # Flatten column names
    quarterly_stats.columns = ['year_quarter', 'avg_unique_bikes', 'std_unique_bikes', 
                             'avg_trips_per_bike', 'std_trips_per_bike']
    
    # Ensure quarterly_stats is sorted
    quarterly_stats = quarterly_stats.sort_values('year_quarter')
    
    # Save utilization metrics to trips_summary.csv
    trips_summary = pd.read_csv(os.path.join(DATA_DIR, TRIPS_FILE))
    trips_summary = trips_summary.merge(quarterly_stats[['year_quarter', 'avg_unique_bikes', 'avg_trips_per_bike']], 
                                      on='year_quarter', how='left')
    trips_summary.to_csv(os.path.join(DATA_DIR, TRIPS_FILE), index=False)
    
    # Save station-level utilization data
    bike_util_stations_df.to_csv(os.path.join(DATA_DIR, STATION_FILE), index=False)
    
    # Plot average unique bikes per station with error bands
    ax1 = plt.gca()
    ax1.plot(range(len(quarterly_stats)), quarterly_stats['avg_unique_bikes'], 
             'b-', label='Avg Unique Bikes per Station', marker='o')
    ax1.fill_between(range(len(quarterly_stats)), 
                    quarterly_stats['avg_unique_bikes'] - quarterly_stats['std_unique_bikes'],
                    quarterly_stats['avg_unique_bikes'] + quarterly_stats['std_unique_bikes'],
                    alpha=0.2, color='blue', label='±1σ Unique Bikes')
    ax1.set_xlabel('Year')
    ax1.set_ylabel('Average Unique Bikes per Station', color='b')
    
    # Plot average trips per bike per station on secondary axis
    ax2 = ax1.twinx()
    ax2.plot(range(len(quarterly_stats)), quarterly_stats['avg_trips_per_bike'], 
             'r-', label='Avg Trips per Bike per Station', marker='o')
    ax2.fill_between(range(len(quarterly_stats)),
                    quarterly_stats['avg_trips_per_bike'] - quarterly_stats['std_trips_per_bike'],
                    quarterly_stats['avg_trips_per_bike'] + quarterly_stats['std_trips_per_bike'],
                    alpha=0.2, color='red', label='±1σ Trips per Bike')
    ax2.set_ylabel('Average Trips per Bike per Station', color='r')
    
    plt.title('Station-level Bike Utilization Trends')
    
    # Update x-axis to show only years
    years = sorted(set([year_quarter.split('Q')[0] for year_quarter in quarterly_stats['year_quarter']]))
    
    # Calculate positions for year labels (at the middle of each year's quarters)
    year_positions = []
    for year in years:
        year_quarters = [i for i, yq in enumerate(quarterly_stats['year_quarter']) if yq.startswith(year)]
        if year_quarters:
            # Place the year label at the middle of its quarters
            year_positions.append(sum(year_quarters) / len(year_quarters))
    
    # Remove tick labels and set year labels
    plt.xticks(year_positions, years)
    ax1.tick_params(axis='x', rotation=0)  # Set rotation to 0 for better readability
    
    # Add legend
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
    
    plt.tight_layout()
    plt.savefig(os.path.join(OUTPUT_DIR, 'station_bike_utilization.png'))
    plt.close()

def main():
    # Create output directory if it doesn't exist
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # Process raw data
    print("Processing raw data...")
    process_raw_data()
    
    # Load and process data
    trips_df, segments_df, station_df, bike_util_system_df, bike_util_stations_df = load_data()
    trips_df = compute_growth_metrics(trips_df)
    trips_df, expansion_stations = detect_anomalies(trips_df, station_df)
    
    # Create visualizations
    create_trends_chart(trips_df)
    create_growth_rate_chart(trips_df)
    create_duration_growth_chart(trips_df)
    create_anomalies_chart(trips_df)
    create_bike_utilization_chart(bike_util_system_df)
    create_station_bike_utilization_chart(bike_util_stations_df)
    
    # Generate presentation
    create_presentation(trips_df, expansion_stations)
    
    print(f"Analysis complete. Presentation saved to {os.path.join(OUTPUT_DIR, OUTPUT_PPTX)}")

if __name__ == "__main__":
    main() 